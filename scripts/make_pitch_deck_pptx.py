from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "pptx"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "clearhire-hackathon-pitch.pptx"

SW, SH = 13.333, 7.5
BG = RGBColor(7, 17, 31)
PANEL = RGBColor(13, 27, 45)
PANEL2 = RGBColor(18, 36, 58)
INK = RGBColor(245, 247, 251)
MUTED = RGBColor(167, 181, 201)
INDIGO = RGBColor(124, 120, 255)
MINT = RGBColor(71, 215, 165)
CORAL = RGBColor(255, 143, 122)
LINE = RGBColor(32, 52, 77)


def box(slide, x, y, w, h, fill=BG, line_color=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color or fill
    shape.line.width = Pt(1)
    return shape


def rect(slide, x, y, w, h, fill):
    return box(slide, x, y, w, h, fill, fill, False)


def txt(slide, value, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return shape


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Inches(width/72))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape


def dot(slide, x, y, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background(); return shape


def chrome(slide, num, label):
    txt(slide, "CLEARHIRE", .58, .25, 1.5, .2, 9, MUTED, True)
    txt(slide, label.upper(), 10.0, .25, 2.75, .2, 9, MUTED, True, PP_ALIGN.RIGHT)
    line(slide, .58, 7.08, 12.75, 7.08)
    txt(slide, f"0{num}", 12.25, 7.17, .5, .16, 9, MUTED, True, PP_ALIGN.RIGHT)


def title(slide, eyebrow, heading, sub):
    txt(slide, eyebrow.upper(), .72, .85, 3.5, .22, 10, MINT, True)
    txt(slide, heading, .72, 1.13, 8.2, 1.15, 30, INK, True)
    txt(slide, sub, .72, 2.35, 7.0, .65, 13, MUTED)


def card(slide, x, y, w, h, heading, body, accent=MINT, metric=None):
    box(slide, x, y, w, h, PANEL, LINE)
    if metric: txt(slide, metric, x+.28, y+.28, w-.56, .5, 30, accent, True)
    txt(slide, heading, x+.28, y+(1.0 if metric else .3), w-.56, .35, 14, INK, True)
    txt(slide, body, x+.28, y+(1.43 if metric else .75), w-.56, h-(1.65 if metric else 1.0), 11, MUTED)


def base(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, SW, SH, BG); return slide


def cover(prs):
    s=base(prs)
    txt(s,"CLEARHIRE",.72,.45,2,.25,11,MINT,True); txt(s,"AI BUILD FEST 2026  /  TRACK 1",9.5,.45,3.1,.25,10,MUTED,True,PP_ALIGN.RIGHT)
    txt(s,"Hiring should not\nbe a memory test.",.72,1.35,7.5,1.65,40,INK,True)
    txt(s,"ClearHire helps growing teams read every application fairly, explain the ranking, and close the loop with every candidate.",.76,3.55,5.6,.75,16,MUTED)
    box(s,8.55,1.3,3.55,3.8,PANEL,LINE)
    txt(s,"THE CONTROL BOUNDARY",8.95,1.75,2.8,.2,10,MINT,True)
    txt(s,"Score merit",8.95,2.35,2.7,.4,27,INK,True); txt(s,"first.",8.95,2.78,2.7,.4,27,INDIGO,True)
    line(s,8.95,3.38,11.85,3.38,LINE,2); txt(s,"Reveal identity",8.95,3.75,2.8,.3,18,INK,True); txt(s,"second.",8.95,4.12,2.8,.3,18,MINT,True); txt(s,"Then let a human decide.",8.95,4.72,2.8,.2,11,MUTED)


def problem(prs):
    s=base(prs); chrome(s,2,"The cost of volume"); title(s,"THE PROBLEM","Recruiters do not have a screening problem. They have a consistency problem.","When applications pile up, the process degrades in predictable ways.")
    card(s,.72,4.0,3.62,2.0,"applications","A realistic SME role can attract more candidates than one person can properly read.",CORAL,"300")
    card(s,4.72,4.0,3.62,2.0,"CVs read","Fatigue turns a structured decision into a first-impression contest.",INDIGO,"40")
    card(s,8.72,4.0,3.62,2.0,"closure","Rejected candidates often receive silence instead of a clear next step.",MINT,"0")


def insight(prs):
    s=base(prs); chrome(s,3,"The product insight"); title(s,"THE INSIGHT","Do not automate the hiring decision. Automate the noise around it.","ClearHire keeps judgment human while making the repeatable work structured, explainable, and fast.")
    for x,n,h,b,a in [(1.0,"01","Blind","Score only job-relevant evidence.",MINT),(4.0,"02","Explain","Show gaps against the rubric.",INDIGO),(7.0,"03","Reveal","Connect the score to a person only after.",CORAL),(10.0,"04","Follow through","Schedule, remind, and close the loop.",MINT)]:
        txt(s,n,x,3.25,2,.2,10,a,True); txt(s,h,x,3.7,2.2,.35,18,INK,True); txt(s,b,x,4.18,2.2,.65,11,MUTED)
    box(s,.72,5.85,11.9,.62,PANEL2); txt(s,"The score is a decision aid, not a decision-maker.",.72,6.04,11.9,.2,14,INK,True,PP_ALIGN.CENTER)


def workflow(prs):
    s=base(prs); chrome(s,4,"One operating system"); title(s,"THE WORKFLOW","From inbox to interview, without losing the human thread.","A single candidate record carries the process from raw CV to final conversation.")
    stages=[("IN","Upload or Gmail","PDF / DOCX\nprivate storage"),("READ","Extract","structured\nprofile"),("RANK","Score blind","weighted rubric\n+ gap analysis"),("MOVE","Schedule","timezone + ICS\n4 reminders"),("CLOSE","Decide","scorecard\nrespectful close")]
    x=.72
    for i,(tag,head,body) in enumerate(stages):
        box(s,x,3.5,2.15,1.9,PANEL,LINE); txt(s,tag,x+.22,3.78,1,.2,9,MINT,True); txt(s,head,x+.22,4.22,1.7,.3,16,INK,True); txt(s,body,x+.22,4.7,1.6,.5,11,MUTED)
        if i<4: txt(s,">",x+2.23,4.22,.25,.3,24,INDIGO,True,PP_ALIGN.CENTER)
        x+=2.42
    txt(s,"One job. One rubric. One explainable shortlist. One place to keep the promise to candidates.",.72,6.05,11.2,.35,14,INK,True)


def blind(prs):
    s=base(prs); chrome(s,5,"The defensible difference"); title(s,"BLIND BY CONSTRUCTION","The model never receives the identity fields.","This is not a blurred table. The blind boundary is enforced where the scoring payload is built.")
    box(s,.72,3.25,5.4,2.65,PANEL,LINE); txt(s,"SENT TO THE MODEL",1.05,3.62,3,.2,10,MINT,True)
    for i,v in enumerate(["skills","experience_years","certifications","tools"]): box(s,1.05,4.02+i*.42,3.3,.3,PANEL2); txt(s,v,1.27,4.1+i*.42,2.8,.14,11,INK)
    box(s,6.82,3.25,5.8,2.65,RGBColor(36,27,42),RGBColor(99,50,76)); txt(s,"NEVER SENT",7.15,3.62,3,.2,10,CORAL,True)
    for i,v in enumerate(["name","email","school","photo / identity"]): txt(s,"x",7.2,4.02+i*.42,.2,.15,14,CORAL,True); txt(s,v,7.55,4.08+i*.42,3,.15,11,INK)


def demo(prs):
    s=base(prs); chrome(s,6,"The money moment"); title(s,"THE DEMO","The reveal is the moment judges remember.","Show the score before the name. Then prove the rest of the workflow is already connected.")
    box(s,.72,3.2,7,2.75,PANEL,LINE); txt(s,"RANKED SHORTLIST",1.05,3.55,3,.2,10,MINT,True)
    for i,(rank,cand,score,col) in enumerate([("#1","Candidate #1","91",MINT),("#2","Candidate #2","88",INDIGO),("#3","Candidate #3","83",CORAL)]):
        y=4.05+i*.58; line(s,1.05,y+.28,7.35,y+.28); dot(s,1.08,y-.02,.22,col); txt(s,rank,1.12,y+.02,.5,.12,8,BG,True,PP_ALIGN.CENTER); txt(s,cand,1.65,y,.0+.1,.2,13,INK,True); txt(s,"blind score · gaps · rationale",1.65,y+.22,3,.14,9,MUTED); txt(s,score,6.75,y+.03,.5,.25,22,INK,True,PP_ALIGN.RIGHT)
    box(s,8.25,3.2,4.38,2.75,PANEL2); txt(s,"REVEAL",8.6,3.55,2,.2,10,MINT,True); txt(s,"The score was locked before anyone connected it to a person.",8.6,4.08,3.5,.9,20,INK,True); line(s,8.6,5.08,12.25,5.08,LINE,2); txt(s,"Then: schedule · remind · compare · close",8.6,5.38,3.5,.25,11,MUTED)


def proof(prs):
    s=base(prs); chrome(s,7,"Proof, not promises"); title(s,"EVIDENCE","We built the trust story into the system.","The winning claim is not that AI is perfect. It is that the workflow is inspectable, bounded, and human-controlled.")
    card(s,.72,3.55,2.82,2.35,"Cross-tenant isolation","Tenant B could not read, update, or delete Tenant A's job data.",MINT,"RLS")
    card(s,3.8,3.55,2.82,2.35,"Reminder cadence","2 days, 1 day, 12 hours, and 2 hours from the confirmed interview.",INDIGO,"4x")
    card(s,6.88,3.55,2.82,2.35,"Blind audit","Every scoring call records only four job-relevant fields.",CORAL,"JSON")
    card(s,9.96,3.55,2.82,2.35,"CV access","Private bucket with short-lived signed download links.",MINT,"5m")


def positioning(prs):
    s=base(prs); chrome(s,8,"Why this can win"); title(s,"POSITIONING","Not another ATS. A fairness-and-follow-through layer for teams that have outgrown spreadsheets.","ClearHire is deliberately narrow: make screening more consistent, then carry the candidate through the moments that usually break.")
    box(s,.72,3.3,5.65,2.55,PANEL,LINE); txt(s,"BROAD ATS",1.08,3.68,2,.2,10,MUTED,True); txt(s,"Many modules.\nHeavy implementation.\nAI as a feature.",1.08,4.2,4.5,1.1,22,INK,True)
    box(s,6.95,3.3,5.65,2.55,PANEL2,INDIGO); txt(s,"CLEARHIRE",7.3,3.68,2,.2,10,MINT,True); txt(s,"Inbox-native intake.\nBlind-score-then-reveal.\nReminders tuned for no-shows.\nHuman decision stays visible.",7.3,4.2,4.5,1.3,20,INK,True)


def roadmap(prs):
    s=base(prs); chrome(s,9,"The business"); title(s,"A CREDIBLE START","Win the first hiring workflow. Expand from there.","The initial customer is an SME recruiter who needs speed and consistency, not another six-month implementation.")
    for i,(tag,head,body,col) in enumerate([("NOW","SME recruiting teams","CV intake, blind ranking, scheduling, closure",MINT),("NEXT","Agencies and hiring teams","shared templates, analytics, candidate exams",INDIGO),("LATER","Hiring infrastructure","audit trails, integrations, team workflows",CORAL)]):
        y=3.6+i*.85; dot(s,1.0,y,.16,col); txt(s,tag,1.35,y+.02,1,.16,10,col,True); txt(s,head,2.45,y,3,.2,16,INK,True); txt(s,body,2.45,y+.27,5,.2,11,MUTED)
    box(s,8.2,3.35,4.3,2.15,PANEL,LINE); txt(s,"THE NORTH STAR",8.55,3.7,3,.2,10,MINT,True); txt(s,"Every candidate gets a fair read, a clear next step, and a respectful close.",8.55,4.2,3.45,.95,21,INK,True)


def close(prs):
    s=base(prs); txt(s,"CLEARHIRE",.72,.45,2,.25,11,MINT,True); txt(s,"The best candidate\nshould not be the\nfirst CV opened.",.72,1.45,7.3,1.7,37,INK,True); txt(s,"Score merit first. Reveal identity second. Keep the human in the loop.",.76,3.75,6.3,.5,18,MUTED)
    box(s,8.65,1.5,3.55,3.55,PANEL,LINE); txt(s,"LIVE DEMO",9.05,1.95,2,.2,10,MINT,True); txt(s,"clearhire-rho",9.05,2.7,3,.35,24,INK,True); txt(s,".vercel.app",9.05,3.12,3,.25,15,INDIGO,True); line(s,9.05,3.7,11.85,3.7,LINE,2); txt(s,"AI BuildFest 2026",9.05,4.08,3,.2,12,MUTED,True); txt(s,"Track 1 · AI for Business",9.05,4.42,3,.2,11,MUTED)


def main():
    prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
    for fn in [cover,problem,insight,workflow,blind,demo,proof,positioning,roadmap,close]: fn(prs)
    prs.save(PPTX); print(PPTX)


if __name__ == "__main__": main()
