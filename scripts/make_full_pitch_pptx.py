from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = Path(r"C:\Users\NEW USER\Downloads\ClearHire-Hackathon-Pitch-Full.pptx")
SW, SH = 13.333, 7.5
BG = RGBColor(7, 17, 31); PANEL = RGBColor(13, 27, 45); PANEL2 = RGBColor(18, 36, 58)
INK = RGBColor(245, 247, 251); MUTED = RGBColor(167, 181, 201); INDIGO = RGBColor(124, 120, 255)
MINT = RGBColor(71, 215, 165); CORAL = RGBColor(255, 143, 122); LINE = RGBColor(32, 52, 77)

def shape(s, kind, x, y, w, h, fill, stroke=None):
    q=s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h)); q.fill.solid(); q.fill.fore_color.rgb=fill; q.line.color.rgb=stroke or fill; return q
def rect(s,x,y,w,h,fill,stroke=None): return shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,fill,stroke)
def text(s,v,x,y,w,h,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT):
    q=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); t=q.text_frame; t.clear(); t.word_wrap=True
    p=t.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=v; r.font.name="Aptos"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; return q
def rule(s,x1,y1,x2,y2,color=LINE): shape(s,MSO_SHAPE.RECTANGLE,x1,y1,x2-x1,.015,color)
def base(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); shape(s,MSO_SHAPE.RECTANGLE,0,0,SW,SH,BG); return s
def chrome(s,n,label):
    text(s,"CLEARHIRE",.58,.25,1.5,.2,9,MUTED,True); text(s,label.upper(),9.3,.25,3.45,.2,9,MUTED,True,PP_ALIGN.RIGHT); rule(s,.58,7.08,12.75,7.08); text(s,f"{n:02d}",12.25,7.17,.5,.16,9,MUTED,True,PP_ALIGN.RIGHT)
def heading(s,eyebrow,title,sub):
    text(s,eyebrow.upper(),.72,.85,4,.2,10,MINT,True); text(s,title,.72,1.13,11.4,1.0,29,INK,True); text(s,sub,.72,2.28,10.2,.55,13,MUTED)
def panel(s,x,y,w,h,head,body,accent=MINT,metric=None):
    rect(s,x,y,w,h,PANEL,LINE)
    if metric: text(s,metric,x+.25,y+.25,w-.5,.42,27,accent,True); yy=y+.88
    else: yy=y+.25
    text(s,head,x+.25,yy,w-.5,.33,13,INK,True); text(s,body,x+.25,yy+.43,w-.5,h-(yy-y)-.62,10.5,MUTED)
def bullets(s,x,y,items,w=5.2):
    for i,(head,body,accent) in enumerate(items):
        yy=y+i*.82; shape(s,MSO_SHAPE.OVAL,x,yy+.04,.12,.12,accent); text(s,head,x+.28,yy,w,.22,13,INK,True); text(s,body,x+.28,yy+.25,w,.32,10.5,MUTED)

def cover(p):
    s=base(p); text(s,"CLEARHIRE",.72,.45,2,.2,11,MINT,True); text(s,"AI BUILDFEST 2026  /  TRACK 1",9.5,.45,3.1,.2,10,MUTED,True,PP_ALIGN.RIGHT)
    text(s,"Hiring should not\nbe a memory test.",.72,1.4,7.5,1.55,40,INK,True); text(s,"An AI recruitment assistant that reads every CV fairly, acts on the shortlist, and never ghosts the candidate.",.76,3.55,6.1,.65,16,MUTED)
    rect(s,8.5,1.3,3.7,3.9,PANEL,LINE); text(s,"THE PROMISE",8.9,1.72,2.8,.2,10,MINT,True); text(s,"Ingest.",8.9,2.25,2.7,.35,25,INK,True); text(s,"Understand.",8.9,2.72,2.7,.35,25,INDIGO,True); text(s,"Act.",8.9,3.19,2.7,.35,25,CORAL,True); rule(s,8.9,3.8,11.75,3.8); text(s,"Keep the final call human.",8.9,4.25,2.8,.3,14,MUTED,True)

def problem(p):
    s=base(p); chrome(s,2,"The cost of volume"); heading(s,"THE PROBLEM","Recruiting breaks in the gaps between tools.","CVs arrive in inboxes. Screening happens under fatigue. Scheduling happens in a dozen messages. Rejections disappear into silence.")
    panel(s,.72,3.45,3.65,2.2,"applications","More candidates than one person can properly read.",CORAL,"300")
    panel(s,4.84,3.45,3.65,2.2,"CVs read","A first-impression contest replaces a consistent rubric.",INDIGO,"40")
    panel(s,8.96,3.45,3.65,2.2,"candidate closure","The people who do not proceed are often left waiting.",MINT,"0")

def insight(p):
    s=base(p); chrome(s,3,"The product insight"); heading(s,"THE INSIGHT","Do not automate the hiring decision. Automate the noise around it.","ClearHire connects the repetitive work around judgment, while keeping judgment visible and human-owned.")
    panel(s,.72,3.35,2.8,2.2,"Blind","Score job-relevant evidence before identity enters the review.",MINT,"01")
    panel(s,3.75,3.35,2.8,2.2,"Copilot","Ask questions and propose actions. Recruiter confirms every side effect.",INDIGO,"02")
    panel(s,6.78,3.35,2.8,2.2,"Follow through","Schedule, email, remind, assess, and score the interview.",CORAL,"03")
    panel(s,9.81,3.35,2.8,2.2,"Human control","The AI recommends. The recruiter decides.",MINT,"04")

def architecture(p):
    s=base(p); chrome(s,4,"One connected system"); heading(s,"THE PLATFORM","One candidate record. Every important handoff connected.","ClearHire turns disconnected recruiting tasks into one inspectable operating flow.")
    stages=[("INBOX","Upload / Gmail"),("AI","Extract / score"),("REVIEW","Shortlist / reveal"),("ACTION","Email / schedule"),("ASSESS","Exam / scorecard"),("CLOSE","Remind / reject")]
    x=.72
    for i,(a,b) in enumerate(stages):
        rect(s,x,3.55,1.8,1.55,PANEL,LINE); text(s,a,x+.2,3.85,1.4,.2,9,MINT,True); text(s,b,x+.2,4.25,1.4,.42,13,INK,True)
        if i<5: text(s,">",x+1.88,4.05,.25,.3,22,INDIGO,True,PP_ALIGN.CENTER)
        x+=2.1
    text(s,"The product is not one AI call. It is the workflow that happens before and after it.",.72,5.95,11.3,.3,15,INK,True)

def gmail(p):
    s=base(p); chrome(s,5,"Inbox-native intake"); heading(s,"GMAIL PULLING","The candidate does not need to learn another application.","Connect Gmail once. ClearHire scans incoming CV attachments, matches them to open roles, and preserves unmatched mail for manual assignment.")
    panel(s,.72,3.3,3.5,2.35,"Connect","OAuth read-only mailbox connection. Refresh tokens are encrypted at rest.",MINT,"01")
    panel(s,4.55,3.3,3.5,2.35,"Match","Subject and body signals route a CV to the best open job.",INDIGO,"02")
    panel(s,8.38,3.3,4.25,2.35,"Never drop","Matched CVs enter the same pipeline as uploads. Unmatched attachments are preserved for review.",CORAL,"03")
    text(s,"Manual pull on a job page or automatic polling through the Cloudflare Worker.",.72,6.05,11,.3,13,MUTED,True)

def ai(p):
    s=base(p); chrome(s,6,"AI pipeline"); heading(s,"AI EXTRACTION + SCORING","Two AI passes. One clear boundary.","Extraction understands the CV. Scoring sees only the structured evidence that matters for the role.")
    panel(s,.72,3.2,3.65,2.55,"Extract","PDF/DOCX text becomes structured skills, experience, certifications, education, and tools.",MINT,"PASS 1")
    panel(s,4.84,3.2,3.65,2.55,"Score blind","The server sends only skills, experience_years, certifications, and tools.",INDIGO,"PASS 2")
    panel(s,8.96,3.2,3.65,2.55,"Explain","Weighted score, missing requirements, hard/nice-to-have gaps, and rationale.",CORAL,"OUTPUT")
    text(s,"The model never does the arithmetic. ClearHire computes the weighted total in code.",.72,6.05,11.2,.3,13,INK,True)

def copilot(p):
    s=base(p); chrome(s,7,"AI Copilot"); heading(s,"COPILOT","Ask the shortlist a question. Turn the answer into a controlled action.","Copilot sees a server-built blind context. It proposes criteria, then the recruiter confirms the actual operation.")
    rect(s,.72,3.15,6.0,2.85,PANEL,LINE); text(s,"ASK AI",1.05,3.52,2,.2,10,MINT,True); rect(s,1.05,4.0,5.35,.45,PANEL2); text(s,"Who is below 60?",1.28,4.13,4.7,.18,12,INK); rect(s,1.05,4.72,5.35,.65,RGBColor(29,47,65)); text(s,"3 candidates match. Review before rejecting.",1.28,4.9,4.8,.2,11,INK)
    bullets(s,7.25,3.45,[("Reject preview","AI proposes a candidate set. Recruiter confirms and can undo.",CORAL),("Exam setup","Create a role-specific exam from the job description.",INDIGO),("CV scan","Find evidence in stored CV text without revealing identity fields.",MINT)],5.1)

def email(p):
    s=base(p); chrome(s,8,"Communication layer"); heading(s,"AI-DRAFTED EMAILS","Every important message starts faster, then stays editable.","ClearHire drafts from templates and context. Recruiters can edit before sending. Every send is logged with the provider result.")
    panel(s,.72,3.35,3.65,2.35,"Interview invite","AI drafts proposed times, role, recruiter signature, and candidate scheduling link.",MINT,"DRAFT")
    panel(s,4.84,3.35,3.65,2.35,"Respectful rejection","Drafts from the gap analysis, avoids hard-gap blame, and never sends without human confirmation.",CORAL,"CLOSE")
    panel(s,8.96,3.35,3.65,2.35,"Delivery trail","Resend sends. ClearHire records recipient, subject, timestamp, and provider message ID.",INDIGO,"LOG")

def scheduling(p):
    s=base(p); chrome(s,9,"Scheduling"); heading(s,"INTERVIEWS + REMINDERS","The candidate picks a time. The system remembers the promise.","Offer up to three slots or book directly. The candidate sees times in their timezone and receives a calendar file.")
    panel(s,.72,3.35,2.8,2.35,"Offer","Three slots or one direct booking.",MINT,"01")
    panel(s,3.75,3.35,2.8,2.35,"Confirm","Candidate chooses from a public personal link.",INDIGO,"02")
    panel(s,6.78,3.35,2.8,2.35,"Calendar","ICS confirmation with location and time.",CORAL,"03")
    panel(s,9.81,3.35,2.8,2.35,"Remember","2d · 1d · 12h · 2h reminder cadence.",MINT,"04")

def exam(p):
    s=base(p); chrome(s,10,"Examination"); heading(s,"EXAMS","From shortlist to evidence, without leaving the job page.","Recruiters choose the candidates, dates, duration, question count, weighting, and tone. ClearHire generates the bank and emails personal links.")
    panel(s,.72,3.25,3.65,2.55,"Recruiter schedules","Set opening and closing dates, choose candidates, and set CV/exam weighting.",MINT,"1")
    panel(s,4.84,3.25,3.65,2.55,"Candidate takes","Personal token link, timed fullscreen attempt, stable question draw, tab-switch strikes.",INDIGO,"2")
    panel(s,8.96,3.25,3.65,2.55,"Shortlist updates","Server grades the attempt. Final score blends CV and exam performance in code.",CORAL,"3")
    text(s,"The exam is not a black box. The recruiter controls the window; the candidate sees the rules before starting.",.72,6.05,11.3,.3,13,INK,True)

def proof(p):
    s=base(p); chrome(s,11,"Proof, not promises"); heading(s,"EVIDENCE","We built inspectability into the product.","The winning claim is not that AI is perfect. It is that the workflow is bounded, testable, and human-controlled.")
    panel(s,.72,3.45,2.82,2.35,"Tenant isolation","RLS tests blocked cross-tenant reads and updates.",MINT,"RLS")
    panel(s,3.8,3.45,2.82,2.35,"Blind audit","Scoring payload contains four job-relevant fields.",INDIGO,"JSON")
    panel(s,6.88,3.45,2.82,2.35,"Private CVs","Private bucket and short-lived signed access.",CORAL,"5m")
    panel(s,9.96,3.45,2.82,2.35,"Live tests","Exam lifecycle, scheduling, auth denial, and isolation tested.",MINT,"E2E")

def close(p):
    s=base(p); text(s,"CLEARHIRE",.72,.45,2,.2,11,MINT,True); text(s,"The best candidate\nshould not be the\nfirst CV opened.",.72,1.4,7.5,1.55,37,INK,True); text(s,"Score merit first. Reveal identity second. Keep the human in the loop.",.76,3.55,6.4,.55,18,MUTED)
    rect(s,8.65,1.45,3.55,3.7,PANEL,LINE); text(s,"LIVE DEMO",9.05,1.9,2,.2,10,MINT,True); text(s,"clearhire-rho",9.05,2.65,3,.35,24,INK,True); text(s,".vercel.app",9.05,3.08,3,.25,15,INDIGO,True); rule(s,9.05,3.68,11.85,3.68); text(s,"AI BuildFest 2026",9.05,4.1,3,.2,12,MUTED,True); text(s,"Track 1 · AI for Business",9.05,4.43,3,.2,11,MUTED)

def main():
    p=Presentation(); p.slide_width=Inches(SW); p.slide_height=Inches(SH)
    for fn in [cover,problem,insight,architecture,gmail,ai,copilot,email,scheduling,exam,proof,close]: fn(p)
    p.save(PPTX); print(PPTX)

if __name__ == "__main__": main()
